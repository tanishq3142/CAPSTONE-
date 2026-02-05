"""
Create PowerPoint Presentation for PolyMix Algorithm Comparison
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path

# Paths
BASE_DIR = Path(__file__).parent
PLOTS_DIR = BASE_DIR / "comparison_plots"
OUTPUT_PATH = BASE_DIR / "PolyMix_Comparison_Presentation.pptx"

# Color scheme
TITLE_COLOR = RGBColor(46, 134, 171)  # Blue
ACCENT_COLOR = RGBColor(233, 79, 55)  # Red


def add_title_slide(prs, title, subtitle):
    """Add a title slide."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_shape.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = TITLE_COLOR
    title_para.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_shape = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1))
    subtitle_frame = subtitle_shape.text_frame
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.text = subtitle
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = RGBColor(100, 100, 100)
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    return slide


def add_section_slide(prs, title):
    """Add a section divider slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Add colored background shape
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(2.5), Inches(10), Inches(2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = TITLE_COLOR
    shape.line.fill.background()
    
    # Title
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(9), Inches(1.2))
    title_frame = title_shape.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.alignment = PP_ALIGN.CENTER
    
    return slide


def add_image_slide(prs, title, image_path, notes=None):
    """Add a slide with an image."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_shape = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9.4), Inches(0.6))
    title_frame = title_shape.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.font.color.rgb = TITLE_COLOR
    title_para.alignment = PP_ALIGN.CENTER
    
    # Image - centered and sized appropriately
    if image_path.exists():
        slide.shapes.add_picture(str(image_path), Inches(0.4), Inches(0.85), width=Inches(9.2))
    
    # Notes if provided
    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes
    
    return slide


def add_text_slide(prs, title, content_lines):
    """Add a slide with text content."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_shape = slide.shapes.add_textbox(Inches(0.3), Inches(0.3), Inches(9.4), Inches(0.8))
    title_frame = title_shape.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = TITLE_COLOR
    title_para.alignment = PP_ALIGN.LEFT
    
    # Content
    content_shape = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.5))
    content_frame = content_shape.text_frame
    content_frame.word_wrap = True
    
    for i, line in enumerate(content_lines):
        if i == 0:
            para = content_frame.paragraphs[0]
        else:
            para = content_frame.add_paragraph()
        para.text = line
        para.font.size = Pt(18)
        para.space_after = Pt(12)
        
        # Highlight key metrics
        if "Winner" in line or "Best" in line:
            para.font.bold = True
    
    return slide


def add_comparison_table_slide(prs, title, headers, data):
    """Add a slide with a comparison table."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_shape = slide.shapes.add_textbox(Inches(0.3), Inches(0.3), Inches(9.4), Inches(0.7))
    title_frame = title_shape.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.font.color.rgb = TITLE_COLOR
    title_para.alignment = PP_ALIGN.CENTER
    
    # Table
    rows = len(data) + 1
    cols = len(headers)
    table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.2), Inches(9), Inches(3)).table
    
    # Style header row
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(14)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = TITLE_COLOR
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    # Data rows
    for row_idx, row_data in enumerate(data):
        for col_idx, value in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(value)
            cell.text_frame.paragraphs[0].font.size = Pt(13)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # Highlight winner column
            if col_idx == len(row_data) - 1 and "Ratio" in str(value):
                cell.text_frame.paragraphs[0].font.bold = True
                if "1.0" in str(value):
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(200, 230, 250)  # Light blue
                else:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(255, 220, 210)  # Light red
    
    return slide


def add_key_findings_slide(prs):
    """Add a key findings summary slide."""
    findings = [
        "📊 Overall Performance:",
        "   • Ratio 2.0 wins in 3/5 training sizes (50, 75, 490)",
        "   • Ratio 1.0 wins in 2/5 training sizes (25, 100)",
        "",
        "🎯 Best Results Achieved:",
        "   • Highest Dice Score: 0.9443 (Ratio 2.0, TS=490)",
        "   • Both ratios achieve >94% Dice with full dataset",
        "",
        "⚡ Convergence Patterns:",
        "   • Smaller datasets (25, 50) require more epochs",
        "   • Larger datasets converge faster (within 15-25 epochs)",
        "",
        "💡 Recommendations:",
        "   • Use Ratio 2.0 for limited data scenarios (50-75 images)",
        "   • Both ratios work similarly well with large datasets"
    ]
    return add_text_slide(prs, "Key Findings & Recommendations", findings)


def main():
    print("Creating PowerPoint presentation...")
    
    # Create presentation with widescreen aspect ratio
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # 1. Title Slide
    add_title_slide(
        prs,
        "PolyMix Algorithm Comparison",
        "Augmentation Ratio 1.0 vs Ratio 2.0\nCapstone Project"
    )
    
    # 2. Experiment Setup Slide
    setup_content = [
        "🔧 Experiment Configuration:",
        "   • Batch Size: 8",
        "   • Learning Rate: 0.0001",
        "   • Image Size: 224 × 224",
        "   • Early Stopping: 10 epochs patience",
        "   • Maximum Epochs: 50",
        "",
        "📁 Training Sizes Tested:",
        "   • 25, 50, 75, 100, 490 images",
        "",
        "🎯 Model Selection Criterion:",
        "   • Best Validation Dice Score"
    ]
    add_text_slide(prs, "Experiment Setup", setup_content)
    
    # 3. Section: Best Performance Comparison
    add_section_slide(prs, "Performance Comparison")
    
    # 4. Best Dice Bar Chart
    add_image_slide(
        prs,
        "Best Validation Dice Score Comparison",
        PLOTS_DIR / "best_dice_comparison.png",
        "Bar chart comparing best dice scores between ratios across training sizes"
    )
    
    # 5. Dice Score Heatmap
    add_image_slide(
        prs,
        "Performance Heatmap Overview",
        PLOTS_DIR / "dice_heatmap.png",
        "Heatmap showing dice scores for both ratios and all training sizes"
    )
    
    # 6. Results Table
    dice_data = [
        ["25", "0.8577", "0.8574", "Ratio 1.0 (+0.03%)"],
        ["50", "0.8691", "0.8827", "Ratio 2.0 (+1.6%)"],
        ["75", "0.8837", "0.8881", "Ratio 2.0 (+0.5%)"],
        ["100", "0.8923", "0.8785", "Ratio 1.0 (+1.4%)"],
        ["490", "0.9421", "0.9443", "Ratio 2.0 (+0.2%)"],
    ]
    add_comparison_table_slide(
        prs,
        "Best Validation Dice Scores",
        ["Training Size", "Ratio 1.0", "Ratio 2.0", "Winner"],
        dice_data
    )
    
    # 7. Section: Training Dynamics
    add_section_slide(prs, "Training Dynamics")
    
    # 8. Convergence Comparison
    add_image_slide(
        prs,
        "Convergence Speed Analysis",
        PLOTS_DIR / "convergence_comparison.png",
        "Number of epochs to reach best dice score"
    )
    
    # 9. Validation Dice Curves
    add_image_slide(
        prs,
        "Validation Dice Across Training Sizes",
        PLOTS_DIR / "val_dice_by_size.png",
        "Epoch-wise validation dice comparison for each training size"
    )
    
    # 10. Size Effect per Ratio
    add_image_slide(
        prs,
        "Training Size Effect per Augmentation Ratio",
        PLOTS_DIR / "size_effect_per_ratio.png",
        "How training size affects performance within each ratio"
    )
    
    # 11. Section: Detailed Training Curves
    add_section_slide(prs, "Detailed Training Curves")
    
    # 12-16. Training Curves for each size
    for ts in [25, 50, 75, 100, 490]:
        add_image_slide(
            prs,
            f"Training Curves - {ts} Images",
            PLOTS_DIR / f"training_curves_ts_{ts}.png",
            f"Detailed training curves for training size {ts}"
        )
    
    # 17. Validation Loss Comparison
    add_image_slide(
        prs,
        "Validation Loss Comparison",
        PLOTS_DIR / "val_loss_by_size.png",
        "Validation loss curves across training sizes"
    )
    
    # 18. Section: Summary
    add_section_slide(prs, "Summary & Conclusions")
    
    # 19. Key Findings
    add_key_findings_slide(prs)
    
    # 20. Conclusion Slide
    conclusion_content = [
        "✅ Conclusions:",
        "",
        "1. Both augmentation ratios achieve strong performance",
        "   with best Dice scores exceeding 94%",
        "",
        "2. Ratio 2.0 shows slight advantage for medium-sized",
        "   datasets (50-75 images)",
        "",
        "3. Ratio 1.0 performs better with very small (25) and",
        "   moderate (100) training sets",
        "",
        "4. Training size has the most significant impact on",
        "   final model performance",
        "",
        "5. Early stopping effectively prevents overfitting",
        "   in all configurations"
    ]
    add_text_slide(prs, "Conclusions", conclusion_content)
    
    # 21. Thank You Slide
    add_title_slide(prs, "Thank You!", "Questions?")
    
    # Save presentation
    prs.save(str(OUTPUT_PATH))
    print(f"\n✅ Presentation saved to: {OUTPUT_PATH}")
    print(f"   Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
